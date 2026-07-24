"""Extract text from all catalog PDFs and PPT files."""
import os
import json
import sys

EXTRACT_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "catalogs_extracted")

results = {}

# Extract PDFs
try:
    import fitz  # PyMuPDF
    for fname in os.listdir(EXTRACT_DIR):
        fpath = os.path.join(EXTRACT_DIR, fname)
        if fname.lower().endswith('.pdf'):
            try:
                doc = fitz.open(fpath)
                text = ""
                for page in doc:
                    text += page.get_text()
                results[fname] = {
                    "type": "pdf",
                    "pages": doc.page_count,
                    "text": text.strip(),
                    "size_bytes": os.path.getsize(fpath)
                }
                doc.close()
            except Exception as e:
                results[fname] = {"type": "pdf", "error": str(e)}
except ImportError:
    print("pymupdf not available", file=sys.stderr)

# Extract PPT
try:
    from pptx import Presentation
    for fname in os.listdir(EXTRACT_DIR):
        fpath = os.path.join(EXTRACT_DIR, fname)
        if fname.lower().endswith(('.ppt', '.pptx')):
            try:
                if fname.lower().endswith('.ppt'):
                    results[fname] = {
                        "type": "ppt",
                        "text": "[Old .ppt format - cannot extract with python-pptx, only .pptx supported]",
                        "size_bytes": os.path.getsize(fpath)
                    }
                else:
                    prs = Presentation(fpath)
                    text = ""
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                for para in shape.text_frame.paragraphs:
                                    text += para.text + "\n"
                    results[fname] = {
                        "type": "pptx",
                        "slides": len(prs.slides),
                        "text": text.strip(),
                        "size_bytes": os.path.getsize(fpath)
                    }
            except Exception as e:
                results[fname] = {"type": "ppt", "error": str(e)}
except ImportError:
    print("python-pptx not available", file=sys.stderr)

# Output - write to file to avoid Windows encoding issues
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "catalog_extracted.json")
with open(output_path, "w", encoding="utf-8") as f:
    json.dump(results, f, indent=2, ensure_ascii=False)
print(f"Extracted {len(results)} files to {output_path}")
